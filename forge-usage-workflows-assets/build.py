# -*- coding: utf-8 -*-
"""forge 사용법·유형별 워크플로우 PPTX 빌더 (python-pptx).
다크 테크 톤 + forge 오렌지 강조. 16:9 와이드. 발표자 노트 포함.
사용: python build.py [출력경로]
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import content as C

# ── 토큰 ─────────────────────────────────────────────────────
BG       = "0F1419"
PANEL    = "1B2531"
PANEL_HI = "243444"
LINE     = "2C3E52"
ACCENT   = "FF7A1A"
ACCENT2  = "4ECDC4"
TEXT     = "EAF0F6"
MUTED    = "8FA0B3"
WHITE    = "FFFFFF"
DANGER   = "E5564A"

FONT_H = "Arial Black"
FONT_B = "Arial"
FONT_M = "Consolas"

EMU_W = 12192000
EMU_H = 6858000
SW = 13.333
SH = 7.5

def hx(h): return RGBColor.from_string(h)

def set_bg(slide, color=BG):
    """풀블리드 배경 사각형."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = hx(color)
    r.line.fill.background()
    r.shadow.inherit = False
    # 맨 뒤로
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Emu(x), Emu(y), Emu(w), Emu(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = hx(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = hx(line)
        s.line.width = Pt(line_w or 1)
    s.shadow.inherit = False
    if radius is not None:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    return s

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=Pt(4), line_spacing=1.0):
    """runs: [(text, size, color, bold, italic, font), ...] 또는 paragraphs list."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, None, None, None, None, None)]]
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size) if size else None
            r.font.color.rgb = hx(color) if color else None
            r.font.bold = bold if bold is not None else None
            r.font.italic = italic if italic is not None else None
            r.font.name = font if font else None
    return tb

def _img_size(path, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw/ih
    w = max_w; h = int(w/ar)
    if h > max_h:
        h = max_h; w = int(h*ar)
    return w, h

def add_image(slide, path, x, y, max_w, max_h, align="center", valign="middle"):
    w, h = _img_size(path, max_w, max_h)
    if align == "center": x = x + (max_w - w)//2
    elif align == "right": x = x + (max_w - w)
    if valign == "middle": y = y + (max_h - h)//2
    elif valign == "bottom": y = y + (max_h - h)
    return slide.shapes.add_picture(str(path), Emu(x), Emu(y), Emu(w), Emu(h))

def add_arrow(slide, x, y, w, h, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = hx(color)
    s.line.fill.background()
    s.shadow.inherit = False
    try: s.adjustments[0] = 0.45; s.adjustments[1] = 0.55
    except Exception: pass
    return s

def chip(slide, x, y, label, color=ACCENT2, w=None):
    """작은 라벨 칩."""
    sz = Pt(11)
    pad = Emu(120000)
    if w is None:
        w = int(len(label) * 100000) + pad*2
    h = Emu(360000)
    r = add_rect(slide, x, y, w, h, fill=color, radius=0.5)
    add_text(slide, x, y+Emu(40000), w, h-Emu(80000),
             [[(label, 11, BG if color not in (PANEL,PANEL_HI) else TEXT, True, False, FONT_B)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w

def footer(slide, left=None, right=None):
    if left:
        add_text(slide, Emu(int(0.6*914400)), Emu(int(7.0*914400)), Emu(int(8*914400)), Emu(int(0.35*914400)),
                 [[(left, 11, MUTED, False, False, FONT_B)]])
    if right:
        add_text(slide, Emu(int(4.7*914400)), Emu(int(7.0*914400)), Emu(int(8.06*914400)), Emu(int(0.35*914400)),
                 [[(right, 11, MUTED, False, False, FONT_B)]], align=PP_ALIGN.RIGHT)

def title_bar(slide, title, kicker=None):
    add_rect(slide, Emu(int(0.6*914400)), Emu(int(0.55*914400)), Emu(int(0.12*914400)), Emu(int(0.62*914400)), fill=ACCENT)
    runs = [[(title, 30, TEXT, True, False, FONT_H)]]
    if kicker:
        runs.append([(kicker, 13, ACCENT, False, True, FONT_B)])
    add_text(slide, Emu(int(0.92*914400)), Emu(int(0.5*914400)), Emu(int(11.8*914400)), Emu(int(0.95*914400)), runs)

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ── kind 렌더러 ───────────────────────────────────────────────
def r_cover(s, d):
    set_bg(s)
    img = C.IMG / d["img"]
    # 상단 60% 이미지 영역을 어둡게 오버레이할 다크 박스 (이미지는 전체 배경처럼)
    add_image(s, img, 0, 0, EMU_W, Emu(int(4.6*914400)), align="center", valign="top")
    # 하단 다크 패널
    add_rect(s, 0, Emu(int(4.55*914400)), EMU_W, Emu(int(2.95*914400)), fill=BG)
    # 상단 그라데이션 대용 어둠 오버레이
    ov = add_rect(s, 0, 0, EMU_W, Emu(int(1.4*914400)), fill=BG)
    ov.fill.transparency = None
    # 킥커
    add_text(s, Emu(int(0.9*914400)), Emu(int(4.95*914400)), Emu(int(11.5*914400)), Emu(int(0.4*914400)),
             [[(d.get("kicker",""), 14, ACCENT2, False, True, FONT_B)]])
    # 제목
    add_text(s, Emu(int(0.85*914400)), Emu(int(5.35*914400)), Emu(int(11.6*914400)), Emu(int(1.3*914400)),
             [[(d["title"], 52, WHITE, True, False, FONT_H)]])
    # 부제
    add_text(s, Emu(int(0.9*914400)), Emu(int(6.45*914400)), Emu(int(11.5*914400)), Emu(int(0.6*914400)),
             [[(d.get("subtitle",""), 19, MUTED, False, False, FONT_B)]])
    # 하단 라인
    add_rect(s, Emu(int(0.9*914400)), Emu(int(7.05*914400)), Emu(int(2.0*914400)), Emu(int(0.05*914400)), fill=ACCENT)

def r_agenda(s, d):
    set_bg(s); title_bar(s, d["title"])
    add_text(s, Emu(int(0.9*914400)), Emu(int(1.65*914400)), Emu(int(11*914400)), Emu(int(0.5*914400)),
             [[(d["lead"], 16, ACCENT2, False, True, FONT_B)]])
    items = d["items"]
    top = int(2.55*914400)
    rh = int(0.95*914400)
    for i, (num, txt) in enumerate(items):
        y = top + i*rh
        add_rect(s, Emu(int(0.9*914400)), Emu(y), Emu(int(11.5*914400)), Emu(int(0.78*914400)), fill=PANEL, radius=0.12)
        add_rect(s, Emu(int(0.9*914400)), Emu(y), Emu(int(0.9*914400)), Emu(int(0.78*914400)), fill=ACCENT, radius=0.12)
        add_text(s, Emu(int(0.9*914400)), Emu(y), Emu(int(0.9*914400)), Emu(int(0.78*914400)),
                 [[(num, 30, BG, True, False, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(int(2.1*914400)), Emu(y), Emu(int(10.0*914400)), Emu(int(0.78*914400)),
                 [[(txt, 18, TEXT, False, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE)

def r_concept(s, d):
    set_bg(s); title_bar(s, d["title"])
    cards = d["cards"]
    n = len(cards)
    cols = 3 if n == 3 else 2
    # 홀수 카드(5+): 위쪽 (n-1)개를 cols열 그리드, 마지막 1개를 아래 풀폭 행으로
    has_full = (n % 2 == 1) and n >= 5
    grid_n = n - 1 if has_full else n
    rows = (grid_n + cols - 1)//cols
    gap = int(0.35*914400)
    mx = int(0.9*914400)
    total_w = EMU_W - mx*2
    cw = (total_w - gap*(cols-1))//cols
    top = int(1.95*914400)
    avail = int(4.9*914400)
    full_h = int(1.3*914400)
    grid_avail = avail - (full_h + gap) if has_full else avail
    if n == 3:
        ch = int(2.9*914400)
        top = int(1.95*914400) + (avail - ch)//2  # 3열 1행 중앙 정렬
    else:
        ch = (grid_avail - gap*(rows-1))//rows
    full_y = top + rows*ch + (rows-1)*gap + gap
    for i, (h, body) in enumerate(cards):
        is_full = has_full and (i == n-1)
        if is_full:
            x = mx; card_w = total_w; card_h = full_h; y = full_y
        else:
            r_ = i//cols; c = i%cols
            x = mx + c*(cw+gap); card_w = cw; card_h = ch
            y = top + r_*(ch+gap)
        add_rect(s, Emu(x), Emu(y), Emu(card_w), Emu(card_h), fill=PANEL, line=LINE, line_w=1, radius=0.08)
        add_rect(s, Emu(x), Emu(y), Emu(int(0.07*914400)), Emu(card_h), fill=ACCENT2)
        add_text(s, Emu(x+int(0.35*914400)), Emu(y+int(card_h*0.16)), Emu(card_w-int(0.6*914400)), Emu(int(0.5*914400)),
                 [[(h, 19, ACCENT, True, False, FONT_H)]])
        body_top = y + int(card_h*0.42)
        body_h = int(card_h*0.50)
        add_text(s, Emu(x+int(0.35*914400)), Emu(body_top), Emu(card_w-int(0.6*914400)), Emu(body_h),
                 [[(body, 14.5 if not is_full else 15, TEXT, False, False, FONT_B)]], line_spacing=1.12)

def r_flow(s, d):
    set_bg(s); title_bar(s, d["title"])
    steps = d["flow"]
    n = len(steps)
    top = int(2.9*914400)
    total_w = EMU_W - int(1.8*914400)
    box_w = int(2.3*914400)
    arrow_w = int(0.5*914400)
    block = box_w + arrow_w
    total = box_w*n + arrow_w*(n-1)
    start_x = (EMU_W - total)//2
    bh = int(1.9*914400)
    for i, (cmd, desc) in enumerate(steps):
        x = start_x + i*block
        is_cmd = cmd.startswith("/") or cmd.startswith("(") or cmd in (" 또는 ", "배포", "목표", "시작", "바로 실행")
        col = ACCENT if (cmd.startswith("/") or is_cmd and cmd.startswith("/")) else PANEL_HI
        # 명령(/로 시작)이면 오렌지, 아니면 패널
        if cmd.startswith("/"):
            fill = ACCENT; tcol = BG
        else:
            fill = PANEL_HI; tcol = TEXT
        add_rect(s, Emu(x), Emu(top), Emu(box_w), Emu(bh), fill=fill, radius=0.1)
        add_text(s, Emu(x), Emu(top+int(0.32*914400)), Emu(box_w), Emu(int(0.6*914400)),
                 [[(cmd, 19, tcol, True, False, FONT_M)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(x), Emu(top+int(1.0*914400)), Emu(box_w), Emu(int(0.6*914400)),
                 [[(desc, 13.5, MUTED if fill==PANEL_HI else BG, False, False, FONT_B)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n-1:
            add_arrow(s, Emu(x+box_w+int(0.05*914400)), Emu(top+int(0.75*914400)), Emu(arrow_w), Emu(int(0.4*914400)), color=ACCENT)
    if d.get("sub"):
        add_text(s, Emu(int(0.9*914400)), Emu(int(5.1*914400)), Emu(int(11.5*914400)), Emu(int(0.5*914400)),
                 [[(d["sub"], 15, ACCENT2, False, True, FONT_B)]], align=PP_ALIGN.CENTER)
    footer(s, d.get("foot_left"), d.get("foot_right"))

def r_scenario(s, d):
    set_bg(s); title_bar(s, d["title"])
    # 상황 배너
    add_rect(s, Emu(int(0.9*914400)), Emu(int(1.7*914400)), Emu(int(11.5*914400)), Emu(int(0.7*914400)), fill=PANEL, line=LINE, line_w=1, radius=0.18)
    add_rect(s, Emu(int(0.9*914400)), Emu(int(1.7*914400)), Emu(int(0.07*914400)), Emu(int(0.7*914400)), fill=ACCENT)
    add_text(s, Emu(int(1.2*914400)), Emu(int(1.7*914400)), Emu(int(11.0*914400)), Emu(int(0.7*914400)),
             [[("상황  ", 13, ACCENT, True, False, FONT_H), (d["situation"], 15.5, TEXT, False, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE)
    # 흐름
    steps = d["flow"]; n = len(steps)
    top = int(2.75*914400)
    box_w = int(2.55*914400)
    arrow_w = int(0.32*914400)
    block = box_w + arrow_w
    total = box_w*n + arrow_w*(n-1)
    start_x = (EMU_W - total)//2
    bh = int(1.25*914400)
    for i, (cmd, desc) in enumerate(steps):
        x = start_x + i*block
        if cmd.startswith("/"):
            fill = ACCENT; tcol = BG
        else:
            fill = PANEL_HI; tcol = TEXT
        add_rect(s, Emu(x), Emu(top), Emu(box_w), Emu(bh), fill=fill, radius=0.12)
        add_text(s, Emu(x), Emu(top+int(0.18*914400)), Emu(box_w), Emu(int(0.5*914400)),
                 [[(cmd, 17, tcol, True, False, FONT_M)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(x), Emu(top+int(0.72*914400)), Emu(box_w), Emu(int(0.45*914400)),
                 [[(desc, 12.5, MUTED if fill==PANEL_HI else BG, False, False, FONT_B)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n-1:
            add_arrow(s, Emu(x+box_w+int(0.02*914400)), Emu(top+int(0.48*914400)), Emu(arrow_w), Emu(int(0.3*914400)), color=ACCENT2)
    # 예시 박스
    ey = int(4.3*914400)
    add_rect(s, Emu(int(0.9*914400)), Emu(ey), Emu(int(11.5*914400)), Emu(int(1.4*914400)), fill=PANEL, line=ACCENT2, line_w=1.25, radius=0.06)
    add_text(s, Emu(int(1.15*914400)), Emu(ey+int(0.16*914400)), Emu(int(2*914400)), Emu(int(0.4*914400)),
             [[("예시", 12, ACCENT2, True, False, FONT_H)]])
    add_text(s, Emu(int(1.15*914400)), Emu(ey+int(0.5*914400)), Emu(int(11.0*914400)), Emu(int(0.85*914400)),
             [[(d["example"], 14, TEXT, False, False, FONT_B)]], line_spacing=1.2)
    # 경고/팁 박스
    wy = int(5.95*914400)
    add_rect(s, Emu(int(0.9*914400)), Emu(wy), Emu(int(11.5*914400)), Emu(int(0.95*914400)), fill=PANEL_HI, line=ACCENT, line_w=1.25, radius=0.1)
    add_text(s, Emu(int(1.15*914400)), Emu(wy), Emu(int(11.0*914400)), Emu(int(0.95*914400)),
             [[("⚠  ", 13, ACCENT, True, False, FONT_H), (d["warn"], 13.5, TEXT, False, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

def r_table(s, d):
    set_bg(s); title_bar(s, d["title"])
    headers = d["headers"]; rows = d["rows"]
    ncol = len(headers); nrow = len(rows)
    top = int(1.75*914400)
    mx = int(0.9*914400)
    total_w = EMU_W - mx*2
    avail_h = int(5.25*914400)
    hh = int(0.55*914400)
    gap = int(0.05*914400)
    # 총 높이 = hh + 0.06 + nrow*rh + (nrow-1)*gap  → rh 역산 (gap 포함)
    rh = min(int(0.72*914400), (avail_h - hh - int(0.06*914400) - gap*max(nrow-1,0))//max(nrow,1))
    # 컬럼 폭: 첫 컬럼 좁게(명령/유틸/상황), 가운데 보통, 마지막 넓게
    if ncol == 3:
        widths = [int(total_w*0.26), int(total_w*0.24), int(total_w*0.50)]
    elif ncol == 2:
        widths = [int(total_w*0.32), int(total_w*0.68)]
    else:
        widths = [total_w//ncol]*ncol
    x0 = mx
    # 헤더
    x = x0
    for c, htext in enumerate(headers):
        add_rect(s, Emu(x), Emu(top), Emu(widths[c]-int(0.04*914400)), Emu(hh), fill=ACCENT, radius=0.06)
        add_text(s, Emu(x+int(0.18*914400)), Emu(top), Emu(widths[c]-int(0.3*914400)), Emu(hh),
                 [[(htext, 13.5, BG, True, False, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
        x += widths[c]
    # 행
    for ridx, row in enumerate(rows):
        y = top + hh + int(0.06*914400) + ridx*(rh+int(0.05*914400))
        x = x0
        zebra = PANEL if ridx%2==0 else PANEL_HI
        for c, val in enumerate(row):
            add_rect(s, Emu(x), Emu(y), Emu(widths[c]-int(0.04*914400)), Emu(rh), fill=zebra, line=LINE, line_w=0.75, radius=0.04)
            is_cmd = c==0 and (val.startswith("/") or val.startswith("fg"))
            col = ACCENT2 if is_cmd else TEXT
            font = FONT_M if is_cmd else FONT_B
            sz = 13 if c==0 else 12.5
            add_text(s, Emu(x+int(0.18*914400)), Emu(y), Emu(widths[c]-int(0.3*914400)), Emu(rh),
                     [[(val, sz, col, is_cmd, False, font)]], anchor=MSO_ANCHOR.MIDDLE)
            x += widths[c]

def r_stat(s, d):
    set_bg(s); title_bar(s, d["title"])
    stats = d["stat"]
    n = len(stats)
    top = int(2.35*914400)
    gap = int(0.35*914400)
    mx = int(0.9*914400)
    cw = (EMU_W - mx*2 - gap*(n-1))//n
    ch = int(3.7*914400)
    for i, (label, cmd, desc, color) in enumerate(stats):
        x = mx + i*(cw+gap)
        add_rect(s, Emu(x), Emu(top), Emu(cw), Emu(ch), fill=PANEL, line=LINE, line_w=1, radius=0.06)
        add_rect(s, Emu(x), Emu(top), Emu(cw), Emu(int(0.12*914400)), fill=color)
        add_text(s, Emu(x+int(0.3*914400)), Emu(top+int(0.35*914400)), Emu(cw-int(0.6*914400)), Emu(int(0.6*914400)),
                 [[(label, 20, color, True, False, FONT_H)]])
        add_text(s, Emu(x+int(0.3*914400)), Emu(top+int(1.05*914400)), Emu(cw-int(0.6*914400)), Emu(int(0.7*914400)),
                 [[(cmd, 21 if len(cmd)<14 else 17, TEXT, True, False, FONT_M)]], anchor=MSO_ANCHOR.TOP)
        add_text(s, Emu(x+int(0.3*914400)), Emu(top+int(1.95*914400)), Emu(cw-int(0.6*914400)), Emu(int(1.6*914400)),
                 [[(desc, 14, MUTED, False, False, FONT_B)]], line_spacing=1.18)
    footer(s, d.get("foot_left"), d.get("foot_right"))

def r_twocol(s, d):
    set_bg(s); title_bar(s, d["title"])
    img = C.IMG / d["img"]
    # 우측 이미지
    ix = int(6.7*914400); iy = int(1.75*914400)
    iw = int(5.8*914400); ih = int(4.9*914400)
    add_rect(s, Emu(ix-int(0.05*914400)), Emu(iy-int(0.05*914400)), Emu(iw+int(0.1*914400)), Emu(ih+int(0.1*914400)), fill=PANEL, radius=0.04)
    add_image(s, img, ix, iy, iw, ih, align="center", valign="middle")
    # 좌측 불릿
    bullets = d["bullets"]
    top = int(2.05*914400)
    bh = int(0.9*914400)
    for i, b in enumerate(bullets):
        y = top + i*bh
        add_rect(s, Emu(int(0.9*914400)), Emu(y+int(0.2*914400)), Emu(int(0.12*914400)), Emu(int(0.12*914400)), fill=ACCENT, radius=0.5)
        add_text(s, Emu(int(1.25*914400)), Emu(y), Emu(int(5.2*914400)), Emu(bh),
                 [[(b, 16, TEXT, False, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    footer(s, d.get("foot_left"), d.get("foot_right"))

def r_cheat(s, d):
    set_bg(s); title_bar(s, d["title"])
    # 좌측 결정표, 우측 작은 이미지
    rows = d["rows"]
    top = int(1.75*914400)
    mx = int(0.9*914400)
    total_w = int(9.6*914400)
    rh = int(0.62*914400)
    hh = int(0.55*914400)
    widths = [int(total_w*0.24), int(total_w*0.52), int(total_w*0.24)]
    heads = ["상황", "명령", "다음"]
    x = mx
    for c, htext in enumerate(heads):
        add_rect(s, Emu(x), Emu(top), Emu(widths[c]-int(0.04*914400)), Emu(hh), fill=ACCENT2, radius=0.06)
        add_text(s, Emu(x+int(0.15*914400)), Emu(top), Emu(widths[c]-int(0.25*914400)), Emu(hh),
                 [[(htext, 12.5, BG, True, False, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
        x += widths[c]
    for ridx, row in enumerate(rows):
        y = top + hh + int(0.05*914400) + ridx*(rh+int(0.04*914400))
        x = mx
        zebra = PANEL if ridx%2==0 else PANEL_HI
        for c, val in enumerate(row):
            add_rect(s, Emu(x), Emu(y), Emu(widths[c]-int(0.04*914400)), Emu(rh), fill=zebra, line=LINE, line_w=0.5, radius=0.04)
            is_cmd = c==1
            col = ACCENT if is_cmd else (ACCENT2 if c==2 else TEXT)
            font = FONT_M if is_cmd else FONT_B
            sz = 12 if is_cmd else 12
            add_text(s, Emu(x+int(0.15*914400)), Emu(y), Emu(widths[c]-int(0.25*914400)), Emu(rh),
                     [[(val, sz, col, is_cmd, False, font)]], anchor=MSO_ANCHOR.MIDDLE)
            x += widths[c]
    # 우측 작은 이미지
    ix = int(10.75*914400); iy = int(1.75*914400)
    iw = int(1.85*914400); ih = int(5.0*914400)
    add_image(s, C.IMG / d["img"], ix, iy, iw, ih, align="center", valign="middle")

def r_section(s, d):
    set_bg(s)
    img = C.IMG / d["img"]
    # 좌측 다크 패널, 우측 이미지 풀블리드
    add_image(s, img, Emu(int(5.0*914400)), 0, Emu(EMU_W-int(5.0*914400)), EMU_H, align="center", valign="middle")
    add_rect(s, 0, 0, Emu(int(5.0*914400)), EMU_H, fill=BG)
    # 경계 라인
    add_rect(s, Emu(int(4.95*914400)), 0, Emu(int(0.06*914400)), EMU_H, fill=ACCENT)
    add_text(s, Emu(int(0.8*914400)), Emu(int(2.5*914400)), Emu(int(4.0*914400)), Emu(int(0.5*914400)),
             [[(d["num"], 16, ACCENT2, True, False, FONT_H)]])
    add_text(s, Emu(int(0.75*914400)), Emu(int(2.95*914400)), Emu(int(4.2*914400)), Emu(int(1.5*914400)),
             [[(d["title"], 34, WHITE, True, False, FONT_H)]], line_spacing=1.0)
    add_rect(s, Emu(int(0.8*914400)), Emu(int(4.4*914400)), Emu(int(1.5*914400)), Emu(int(0.06*914400)), fill=ACCENT)

def r_sources(s, d):
    set_bg(s); title_bar(s, d["title"])
    rows = d["rows"]
    top = int(2.0*914400)
    rh = int(0.85*914400)
    gap = int(0.18*914400)
    for i, (name, desc) in enumerate(rows):
        y = top + i*(rh+gap)
        add_rect(s, Emu(int(0.9*914400)), Emu(y), Emu(int(11.5*914400)), Emu(rh), fill=PANEL, line=LINE, line_w=1, radius=0.08)
        add_rect(s, Emu(int(0.9*914400)), Emu(y), Emu(int(0.07*914400)), Emu(rh), fill=ACCENT2)
        add_text(s, Emu(int(1.25*914400)), Emu(y), Emu(int(3.6*914400)), Emu(rh),
                 [[(name, 17, ACCENT2, True, False, FONT_M)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(int(4.9*914400)), Emu(y), Emu(int(7.3*914400)), Emu(rh),
                 [[(desc, 13.5, TEXT, False, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, d.get("foot_left"))

RENDER = {
    "cover": r_cover, "agenda": r_agenda, "concept": r_concept, "flow": r_flow,
    "scenario": r_scenario, "table": r_table, "stat": r_stat, "twocol": r_twocol,
    "cheat": r_cheat, "section": r_section, "sources": r_sources,
}

def build(out_path):
    prs = Presentation()
    prs.slide_width = Emu(EMU_W)
    prs.slide_height = Emu(EMU_H)
    blank = prs.slide_layouts[6]
    for d in C.SLIDES:
        s = prs.slides.add_slide(blank)
        RENDER[d["kind"]](s, d)
        set_notes(s, d.get("notes", ""))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(C.SLIDES)

if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("/Users/gyuha/workspace/forge-learn/forge-사용법-워크플로우.pptx")
    n = build(out)
    print(f"Built {n} slides -> {out}")
